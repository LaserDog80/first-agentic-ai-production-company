"""Convert PitchDeck JSON to PowerPoint (.pptx) files."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# Slide dimensions (standard 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)

# Font sizes
TITLE_SIZE = Pt(32)
SUBTITLE_SIZE = Pt(18)
HEADING_SIZE = Pt(28)
BODY_SIZE = Pt(16)
SMALL_SIZE = Pt(12)
TABLE_SIZE = Pt(10)

# Layout margins
LEFT_MARGIN = Inches(0.8)
TOP_MARGIN = Inches(1.2)
CONTENT_WIDTH = Inches(11.7)
CONTENT_HEIGHT = Inches(5.5)

MAX_CELL_CHARS = 150


def export_pitch_deck(pitch_deck: dict, output_path: str) -> Path:
    """Convert a PitchDeck dict to a .pptx file.

    Args:
        pitch_deck: The PitchDeck dict (same structure as pitch_deck.json).
        output_path: File path for the output .pptx file.

    Returns:
        Path to the written .pptx file.
    """
    path = Path(output_path)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Title
    _add_title_slide(prs, pitch_deck.get("title_page", {}))
    # Slide 2: Logline
    _add_logline_slide(prs, pitch_deck.get("logline", ""))
    # Slide 3: Format & Tone
    _add_format_slide(prs, pitch_deck.get("format_and_tone", {}))
    # Slide 4: Target Audience
    _add_audience_slide(prs, pitch_deck.get("target_audience", ""))
    # Slide 5: Competitive Landscape
    _add_competitors_slide(prs, pitch_deck.get("competitive_landscape", []))
    # Slide 6: Key Characters
    _add_characters_slide(prs, pitch_deck.get("key_characters", []))
    # Slides 7-9: Episode breakdown
    episode = pitch_deck.get("episode_breakdown", {})
    _add_narrative_arc_slide(prs, episode)
    _add_sequences_slide(prs, episode.get("key_sequences", []))
    _add_visual_approach_slide(prs, episode)
    # Slide 10: Feasibility
    _add_feasibility_slide(prs, pitch_deck.get("feasibility_summary", {}))
    # Slide 11: Why Now
    _add_why_now_slide(prs, pitch_deck.get("why_now", ""))
    # Slide 12: Review Notes & Concerns
    _add_review_slide(
        prs,
        pitch_deck.get("sp_review_notes", ""),
        pitch_deck.get("unresolved_concerns", []),
    )

    prs.save(str(path))
    return path


def _add_text_box(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_size: Pt = BODY_SIZE,
    bold: bool = False,
    color: RGBColor = DARK_GREY,
    alignment: int = PP_ALIGN.LEFT,
) -> None:
    """Add a text box to a slide."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment


def _add_slide_heading(slide, text: str) -> None:
    """Add a heading at the top of a content slide."""
    _add_text_box(
        slide,
        left=LEFT_MARGIN,
        top=Inches(0.4),
        width=CONTENT_WIDTH,
        height=Inches(0.7),
        text=text,
        font_size=HEADING_SIZE,
        bold=True,
        color=DARK_BLUE,
    )


def _add_logline_slide(prs: Presentation, logline: str) -> None:
    """Slide 2: Logline — prominent centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_heading(slide, "Logline")
    _add_text_box(
        slide,
        left=LEFT_MARGIN,
        top=Inches(2.0),
        width=CONTENT_WIDTH,
        height=Inches(4.0),
        text=logline,
        font_size=SUBTITLE_SIZE,
        color=DARK_GREY,
        alignment=PP_ALIGN.CENTER,
    )


def _add_format_slide(prs: Presentation, format_and_tone: dict) -> None:
    """Slide 3: Format & Tone — bullet list."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_heading(slide, "Format & Tone")

    items = [
        f"Series Length: {format_and_tone.get('series_length', 'N/A')}",
        f"Genre: {format_and_tone.get('genre', 'N/A')}",
        f"Tone: {format_and_tone.get('tone', 'N/A')}",
    ]

    txbox = slide.shapes.add_textbox(
        LEFT_MARGIN, TOP_MARGIN, CONTENT_WIDTH, CONTENT_HEIGHT
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = BODY_SIZE
        p.font.color.rgb = DARK_GREY
        p.space_after = Pt(12)


def _add_audience_slide(prs: Presentation, target_audience: str) -> None:
    """Slide 4: Target Audience — text block."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_heading(slide, "Target Audience")
    _add_text_box(
        slide,
        left=LEFT_MARGIN,
        top=TOP_MARGIN,
        width=CONTENT_WIDTH,
        height=CONTENT_HEIGHT,
        text=target_audience,
        font_size=BODY_SIZE,
    )


def _truncate(text: str, max_chars: int = MAX_CELL_CHARS) -> str:
    """Truncate text with ellipsis if over max length."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars - 3] + "..."


def _add_table_slide(
    prs: Presentation,
    heading: str,
    headers: list[str],
    rows: list[list[str]],
) -> None:
    """Add a slide with a heading and a styled table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_heading(slide, heading)

    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        LEFT_MARGIN, TOP_MARGIN, CONTENT_WIDTH, CONTENT_HEIGHT,
    )
    table = table_shape.table

    # Header row
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = TABLE_SIZE
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = _truncate(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = TABLE_SIZE
                paragraph.font.color.rgb = BLACK
            # Alternating row color
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GREY


def _add_competitors_slide(prs: Presentation, competitors: list[dict]) -> None:
    """Slide 5: Competitive Landscape — table."""
    rows = [
        [
            c.get("title", ""),
            c.get("broadcaster", ""),
            c.get("year", ""),
            c.get("relevance", ""),
        ]
        for c in competitors
    ]
    _add_table_slide(
        prs, "Competitive Landscape",
        ["Title", "Broadcaster", "Year", "Relevance"],
        rows,
    )


def _add_characters_slide(prs: Presentation, characters: list[dict]) -> None:
    """Slide 6: Key Characters — table."""
    rows = [
        [
            c.get("name", ""),
            c.get("role", ""),
            c.get("access_notes", ""),
            c.get("story_angle", ""),
        ]
        for c in characters
    ]
    _add_table_slide(
        prs, "Key Characters",
        ["Name", "Role", "Access", "Story Angle"],
        rows,
    )


def _add_narrative_arc_slide(prs: Presentation, episode: dict) -> None:
    """Slide 7: Narrative Arc — four labelled text blocks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ep_title = episode.get("episode_title", "Episode Breakdown")
    _add_slide_heading(slide, f"Narrative Arc: {ep_title}")

    arc = episode.get("narrative_arc", {})
    sections = [
        ("Opening", arc.get("opening", "")),
        ("Development", arc.get("development", "")),
        ("Climax", arc.get("climax", "")),
        ("Resolution", arc.get("resolution", "")),
    ]

    txbox = slide.shapes.add_textbox(
        LEFT_MARGIN, TOP_MARGIN, CONTENT_WIDTH, CONTENT_HEIGHT
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, (label, content) in enumerate(sections):
        # Label paragraph (bold)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = label
        p.font.size = SMALL_SIZE
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.space_before = Pt(8) if i > 0 else Pt(0)

        # Content paragraph
        p2 = tf.add_paragraph()
        p2.text = _truncate(content, 300)
        p2.font.size = SMALL_SIZE
        p2.font.color.rgb = DARK_GREY
        p2.space_after = Pt(6)


def _add_sequences_slide(prs: Presentation, sequences: list[dict]) -> None:
    """Slide 8: Key Sequences — table."""
    rows = [
        [
            s.get("name", ""),
            s.get("description", ""),
            s.get("visual_style", ""),
            str(s.get("duration_mins", "")),
        ]
        for s in sequences
    ]
    _add_table_slide(
        prs, "Key Sequences",
        ["Name", "Description", "Visual Style", "Duration (min)"],
        rows,
    )


def _add_visual_approach_slide(prs: Presentation, episode: dict) -> None:
    """Slide 9: Visual Approach — tone + visual approach text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_heading(slide, "Visual Approach")

    tone = episode.get("overall_tone", "")
    approach = episode.get("visual_approach", "")

    txbox = slide.shapes.add_textbox(
        LEFT_MARGIN, TOP_MARGIN, CONTENT_WIDTH, CONTENT_HEIGHT
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    # Tone
    p = tf.paragraphs[0]
    p.text = "Overall Tone"
    p.font.size = BODY_SIZE
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    p2 = tf.add_paragraph()
    p2.text = tone
    p2.font.size = BODY_SIZE
    p2.font.color.rgb = DARK_GREY
    p2.space_after = Pt(18)

    # Visual approach
    p3 = tf.add_paragraph()
    p3.text = "Visual Approach"
    p3.font.size = BODY_SIZE
    p3.font.bold = True
    p3.font.color.rgb = DARK_BLUE

    p4 = tf.add_paragraph()
    p4.text = approach
    p4.font.size = BODY_SIZE
    p4.font.color.rgb = DARK_GREY


def _add_feasibility_slide(prs: Presentation, feasibility: dict) -> None:
    """Slide 10: Feasibility — rating, budget, days, risks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_heading(slide, "Feasibility")

    rating = feasibility.get("feasibility_rating", "N/A").upper()
    budget = feasibility.get("budget_bracket", {})
    low = budget.get("low", "?")
    high = budget.get("high", "?")
    currency = budget.get("currency", "")
    days = feasibility.get("shooting_days", "N/A")
    risks = feasibility.get("key_risks", [])

    txbox = slide.shapes.add_textbox(
        LEFT_MARGIN, TOP_MARGIN, CONTENT_WIDTH, CONTENT_HEIGHT
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    items = [
        f"Rating: {rating}",
        f"Budget: {currency} {low:,} - {high:,}" if isinstance(low, int) else f"Budget: {currency} {low} - {high}",
        f"Shooting Days: {days}",
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = BODY_SIZE
        p.font.bold = True
        p.font.color.rgb = DARK_GREY
        p.space_after = Pt(8)

    # Risks
    if risks:
        p_label = tf.add_paragraph()
        p_label.text = "Key Risks:"
        p_label.font.size = BODY_SIZE
        p_label.font.bold = True
        p_label.font.color.rgb = DARK_BLUE
        p_label.space_before = Pt(12)

        for risk in risks:
            p_risk = tf.add_paragraph()
            p_risk.text = f"  \u2022  {_truncate(risk)}"
            p_risk.font.size = SMALL_SIZE
            p_risk.font.color.rgb = DARK_GREY


def _add_review_slide(
    prs: Presentation, review_notes: str, concerns: list[str]
) -> None:
    """Slide 12: SP Review Notes & Unresolved Concerns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_heading(slide, "Review Notes & Concerns")

    txbox = slide.shapes.add_textbox(
        LEFT_MARGIN, TOP_MARGIN, CONTENT_WIDTH, CONTENT_HEIGHT
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    # Review notes
    p = tf.paragraphs[0]
    p.text = review_notes
    p.font.size = SMALL_SIZE
    p.font.color.rgb = DARK_GREY
    p.space_after = Pt(14)

    # Concerns
    if concerns:
        p_label = tf.add_paragraph()
        p_label.text = "Unresolved Concerns:"
        p_label.font.size = BODY_SIZE
        p_label.font.bold = True
        p_label.font.color.rgb = DARK_BLUE
        p_label.space_before = Pt(12)

        for concern in concerns:
            p_c = tf.add_paragraph()
            p_c.text = f"  \u2022  {_truncate(concern)}"
            p_c.font.size = SMALL_SIZE
            p_c.font.color.rgb = DARK_GREY


def _add_why_now_slide(prs: Presentation, why_now: str) -> None:
    """Slide 11: Why Now — text block."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_heading(slide, "Why Now")
    _add_text_box(
        slide,
        left=LEFT_MARGIN,
        top=TOP_MARGIN,
        width=CONTENT_WIDTH,
        height=CONTENT_HEIGHT,
        text=why_now,
        font_size=BODY_SIZE,
    )


def _add_title_slide(prs: Presentation, title_page: dict) -> None:
    """Slide 1: Title page with working title, genre, format, broadcaster."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    title = title_page.get("working_title", "Untitled")
    genre = title_page.get("genre", "")
    fmt = title_page.get("format", "")
    broadcaster = title_page.get("target_broadcaster", "")

    # Title — centered, large
    _add_text_box(
        slide,
        left=LEFT_MARGIN,
        top=Inches(2.0),
        width=CONTENT_WIDTH,
        height=Inches(1.5),
        text=title,
        font_size=TITLE_SIZE,
        bold=True,
        color=DARK_BLUE,
        alignment=PP_ALIGN.CENTER,
    )

    # Subtitle line: genre | format | broadcaster
    subtitle_parts = [p for p in [genre, fmt, broadcaster] if p]
    subtitle = "  |  ".join(subtitle_parts)
    _add_text_box(
        slide,
        left=LEFT_MARGIN,
        top=Inches(3.8),
        width=CONTENT_WIDTH,
        height=Inches(0.8),
        text=subtitle,
        font_size=SUBTITLE_SIZE,
        color=DARK_GREY,
        alignment=PP_ALIGN.CENTER,
    )
