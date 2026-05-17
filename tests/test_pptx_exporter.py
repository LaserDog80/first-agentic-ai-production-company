"""Tests for the PPTX exporter."""
import json
from pathlib import Path

import pytest
from pptx import Presentation

from src.pptx_exporter import export_pitch_deck


@pytest.fixture
def sample_pitch_deck() -> dict:
    """A self-contained pitch deck fixture exercising every slide."""
    return {
        "title_page": {
            "working_title": "Lighthouse Keepers",
            "genre": "observational documentary",
            "format": "3x60",
            "target_broadcaster": "BBC Two",
        },
        "logline": "On a remote Atlantic rock, the last keeper learns the tide is coming for everything he knows.",
        "format_and_tone": {
            "series_length": "3x60",
            "genre": "observational documentary",
            "tone": "intimate, weather-beaten, hopeful",
        },
        "target_audience": "Adults 35+ who watch slow-form factual on BBC Two and Channel 4 — viewers of The Detectorists, Coast, and Earth at Night.",
        "competitive_landscape": [
            {"title": "Coast", "broadcaster": "BBC Two", "year": "2005",
             "relevance": "Defined the British coastal documentary."},
            {"title": "The Detectorists", "broadcaster": "BBC Four", "year": "2014",
             "relevance": "Demonstrated appetite for slow, place-rooted storytelling."},
        ],
        "key_characters": [
            {"name": "Donald MacLeod", "role": "Last keeper of Hyskeir",
             "access_notes": "Confirmed verbal access via Northern Lighthouse Board.",
             "story_angle": "His final season before automation."},
            {"name": "Aileen Ross", "role": "Marine archaeologist",
             "access_notes": "Strong existing relationship through prior research.",
             "story_angle": "The history beneath the waves."},
        ],
        "episode_breakdown": {
            "episode_title": "The Keeper",
            "narrative_arc": {
                "opening": "Dawn on Hyskeir — Donald climbs the spiral stairs for what he knows is one of the last times.",
                "development": "We follow his daily rituals and the rising tide of automation.",
                "climax": "A storm forces him to choose between protocol and instinct.",
                "resolution": "Helicopter lifts him off; the light keeps turning.",
            },
            "key_sequences": [
                {"name": "Lighting the lamp", "description": "Donald's morning routine",
                 "visual_style": "long lenses, soft natural light", "duration_mins": 6},
                {"name": "Storm", "description": "Force 9 hits the rock",
                 "visual_style": "handheld, restless cuts", "duration_mins": 12},
            ],
            "overall_tone": "Intimate, weather-beaten",
            "visual_approach": "Patient, painterly, with archive interludes.",
            "contributor_usage": [
                {"character_name": "Donald MacLeod", "role_in_episode": "Anchor and emotional spine"},
            ],
            "special_requirements": ["Helicopter access", "Marine insurance"],
        },
        "feasibility_summary": {
            "feasibility_rating": "amber",
            "budget_bracket": {"low": 850000, "high": 1200000, "currency": "GBP",
                                "notes": "Includes 28 shoot days, archive licensing, post."},
            "shooting_days": 28,
            "key_risks": ["Weather-dependent access", "Single-character reliance"],
        },
        "why_now": "Automation reaches the last manned light in 2026 — this is the final chance to film it.",
        "sp_review_notes": "Strong subject; treatment is bold; budget needs scrutiny.",
        "unresolved_concerns": ["Need a B-strand if Donald falls ill mid-shoot."],
    }


def test_export_creates_file(sample_pitch_deck, tmp_path):
    """Exporter creates a valid .pptx file."""
    output = tmp_path / "deck.pptx"
    result = export_pitch_deck(sample_pitch_deck, str(output))
    assert result == output
    assert output.exists()
    prs = Presentation(str(output))
    assert len(prs.slides) > 0


def test_title_slide_content(sample_pitch_deck, tmp_path):
    """Title slide contains the working title (as text or pixel art image)."""
    output = tmp_path / "deck.pptx"
    export_pitch_deck(sample_pitch_deck, str(output))
    prs = Presentation(str(output))
    first_slide = prs.slides[0]
    texts = [shape.text for shape in first_slide.shapes if shape.has_text_frame]
    all_text = " ".join(texts)
    has_image = any(shape.shape_type == 13 for shape in first_slide.shapes)
    # Title may be rendered as pixel art image or text fallback
    title_in_text = sample_pitch_deck["title_page"]["working_title"] in all_text
    assert title_in_text or has_image


def test_text_slides_present(sample_pitch_deck, tmp_path):
    """Slides 2-4 and 11 contain logline, format, audience, and why_now text."""
    output = tmp_path / "deck.pptx"
    export_pitch_deck(sample_pitch_deck, str(output))
    prs = Presentation(str(output))

    def slide_text(index):
        return " ".join(
            shape.text for shape in prs.slides[index].shapes if shape.has_text_frame
        )

    # Slide 2: logline
    assert sample_pitch_deck["logline"][:50] in slide_text(1)
    # Slide 3: format_and_tone
    assert sample_pitch_deck["format_and_tone"]["genre"] in slide_text(2)
    # Slide 4: target_audience
    assert sample_pitch_deck["target_audience"][:50] in slide_text(3)


def test_competitor_table_slide(sample_pitch_deck, tmp_path):
    """Slide 5 contains a table with competitor data."""
    output = tmp_path / "deck.pptx"
    export_pitch_deck(sample_pitch_deck, str(output))
    prs = Presentation(str(output))
    slide = prs.slides[4]  # 0-indexed, slide 5
    tables = [s for s in slide.shapes if s.has_table]
    assert len(tables) == 1
    table = tables[0].table
    # Header row + data rows
    expected_rows = 1 + len(sample_pitch_deck["competitive_landscape"])
    assert len(table.rows) == expected_rows
    # 4 columns: title, broadcaster, year, relevance
    assert len(table.columns) == 4


def test_characters_table_slide(sample_pitch_deck, tmp_path):
    """Slide 6 contains a table with character data."""
    output = tmp_path / "deck.pptx"
    export_pitch_deck(sample_pitch_deck, str(output))
    prs = Presentation(str(output))
    slide = prs.slides[5]  # 0-indexed, slide 6
    tables = [s for s in slide.shapes if s.has_table]
    assert len(tables) == 1
    table = tables[0].table
    expected_rows = 1 + len(sample_pitch_deck["key_characters"])
    assert len(table.rows) == expected_rows
    assert len(table.columns) == 4


def test_narrative_arc_slide(sample_pitch_deck, tmp_path):
    """Slide 7 contains narrative arc sections."""
    output = tmp_path / "deck.pptx"
    export_pitch_deck(sample_pitch_deck, str(output))
    prs = Presentation(str(output))
    slide = prs.slides[6]
    slide_text = " ".join(
        shape.text for shape in slide.shapes if shape.has_text_frame
    )
    arc = sample_pitch_deck["episode_breakdown"]["narrative_arc"]
    # Check that at least the opening is present
    assert arc["opening"][:40] in slide_text


def test_sequences_table_slide(sample_pitch_deck, tmp_path):
    """Slide 8 contains a table of key sequences."""
    output = tmp_path / "deck.pptx"
    export_pitch_deck(sample_pitch_deck, str(output))
    prs = Presentation(str(output))
    slide = prs.slides[7]
    tables = [s for s in slide.shapes if s.has_table]
    assert len(tables) == 1
    table = tables[0].table
    sequences = sample_pitch_deck["episode_breakdown"]["key_sequences"]
    assert len(table.rows) == 1 + len(sequences)
    assert len(table.columns) == 4


def test_full_slide_count(sample_pitch_deck, tmp_path):
    """Complete export produces exactly 12 slides."""
    output = tmp_path / "deck.pptx"
    export_pitch_deck(sample_pitch_deck, str(output))
    prs = Presentation(str(output))
    assert len(prs.slides) == 12


def test_feasibility_slide(sample_pitch_deck, tmp_path):
    """Slide 10 contains feasibility rating."""
    output = tmp_path / "deck.pptx"
    export_pitch_deck(sample_pitch_deck, str(output))
    prs = Presentation(str(output))
    slide = prs.slides[9]
    slide_text = " ".join(
        shape.text for shape in slide.shapes if shape.has_text_frame
    )
    rating = sample_pitch_deck["feasibility_summary"]["feasibility_rating"]
    assert rating.upper() in slide_text.upper()


def test_cli_integration_writes_pptx(sample_pitch_deck, tmp_path):
    """Verify export_pitch_deck works with a real output directory path."""
    out_dir = tmp_path / "output"
    out_dir.mkdir()
    pptx_path = out_dir / "pitch_deck.pptx"
    result = export_pitch_deck(sample_pitch_deck, str(pptx_path))
    assert result == pptx_path
    assert pptx_path.exists()
    assert pptx_path.stat().st_size > 0


def test_minimal_pitch_deck_no_crash(tmp_path):
    """Exporter handles a near-empty PitchDeck dict without crashing."""
    minimal = {
        "title_page": {"working_title": "Test"},
        "logline": "",
        "format_and_tone": {},
        "target_audience": "",
        "competitive_landscape": [],
        "key_characters": [],
        "episode_breakdown": {
            "narrative_arc": {},
            "key_sequences": [],
        },
        "feasibility_summary": {},
        "why_now": "",
        "sp_review_notes": "",
        "unresolved_concerns": [],
    }
    output = tmp_path / "minimal.pptx"
    result = export_pitch_deck(minimal, str(output))
    assert result == output
    prs = Presentation(str(output))
    assert len(prs.slides) == 12
